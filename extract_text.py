import os
from pathlib import Path
import docx
import pptx
import fitz  # PyMuPDF
import json

SOURCE_DIR = "/Users/vasanthagullapalli/Downloads/Pitt_2023-25"
OUTPUT_FILE = "data/raw_text.jsonl"
SUPPORTED_EXTS = {"pdf", "docx", "pptx", "txt"}

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_pptx(path):
    prs = pptx.Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_all_text():
    os.makedirs("data", exist_ok=True)
    with open(OUTPUT_FILE, "w") as outfile:
        for root, _, files in os.walk(SOURCE_DIR):
            for file in files:
                ext = file.lower().split('.')[-1]
                path = os.path.join(root, file)
                content = ""
                if ext == "pdf":
                    content = extract_text_from_pdf(path)
                elif ext == "docx":
                    content = extract_text_from_docx(path)
                elif ext == "pptx":
                    content = extract_text_from_pptx(path)
                elif ext == "txt":
                    try:
                        content = Path(path).read_text()
                    except Exception as e:
                        print(f"[TXT ERROR] Failed to extract {path}: {e}")
                        content = ""
                else:
                    continue  # Skip unsupported formats

                if content.strip():  # Only write non-empty content
                    json.dump({"file": path, "text": content}, outfile)
                    outfile.write("\n")

if __name__ == "__main__":
    extract_all_text()
